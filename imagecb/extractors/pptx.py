"""Extract images plus slide context from .pptx files."""

from __future__ import annotations

import io
import logging
from datetime import datetime
from pathlib import Path
from typing import Iterator

from PIL import Image

from .pptx_text import slide_body_text, slide_notes, slide_title
from .types import ExtractedImage, Provenance

logger = logging.getLogger(__name__)

_MAX_BODY_CHARS = 1500


def _truncate(text: str, limit: int) -> str:
    text = (text or "").strip()
    if len(text) <= limit:
        return text
    return text[: limit - 3].rstrip() + "..."


def _iter_picture_shapes(shapes):
    """Recursively yield picture shapes, descending into group shapes."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            yield shape
        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_picture_shapes(shape.shapes)


def extract(path: Path) -> Iterator[ExtractedImage]:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        logger.error("python-pptx not installed: %s", exc)
        return

    try:
        prs = Presentation(str(path))
    except Exception as exc:  # noqa: BLE001
        logger.warning("Failed to open pptx %s: %s", path, exc)
        return

    core = prs.core_properties
    author = core.author or core.last_modified_by or None
    created_at = core.created
    modified_at = core.modified
    if not modified_at:
        modified_at = datetime.fromtimestamp(path.stat().st_mtime)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        title = slide_title(slide)
        notes = slide_notes(slide)
        body = _truncate(slide_body_text(slide, title=title), _MAX_BODY_CHARS)
        for shape in _iter_picture_shapes(slide.shapes):
            try:
                blob = shape.image.blob
                img = Image.open(io.BytesIO(blob))
                img.load()
            except Exception as exc:  # noqa: BLE001
                logger.debug("Could not load image on slide %s of %s: %s", slide_idx, path, exc)
                continue
            provenance = Provenance(
                source_file=str(path),
                source_type="pptx",
                source_modified_at=modified_at,
                source_created_at=created_at,
                author=author,
                slide_index=slide_idx,
                slide_title=title,
                slide_notes=notes,
                extra={"slide_body_text": body} if body else {},
            )
            yield ExtractedImage(image=img, provenance=provenance)
